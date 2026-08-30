import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. 16:9 와이드스크린 프레젠테이션 생성 (13.333" x 7.5" / 1920x1080)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Pretendard 단일 폰트 & 엄격한 디자인 시스템 색상 토큰
FONT_NAME = "Pretendard"

# Brand & Accent Colors
COLOR_PRIMARY = RGBColor(0, 102, 204)       # Primary Action Blue (#0066CC)
COLOR_FOCUS_BLUE = RGBColor(41, 151, 255)   # Electric Focus Blue (#2997FF)
COLOR_EMERALD = RGBColor(16, 185, 129)      # Emerald Accent (#10B981)
COLOR_AMBER = RGBColor(245, 158, 11)        # Amber Warning (#F59E0B)
COLOR_PURPLE = RGBColor(139, 92, 246)       # Purple Flow Node (#8B5CF6)
COLOR_ROSE = RGBColor(244, 63, 94)          # Rose Accent (#F43F5E)

# Surface Colors
COLOR_CANVAS_LIGHT = RGBColor(248, 250, 252)      # Canvas Light (#F8FAFC)
COLOR_CANVAS_PARCHMENT = RGBColor(245, 245, 247)  # Canvas Parchment (#F5F5F7)
COLOR_CANVAS_DARK = RGBColor(30, 41, 59)          # Canvas Slate Navy (#1E293B)
COLOR_CARD_WHITE = RGBColor(255, 255, 255)        # Surface Card White (#FFFFFF)
COLOR_CARD_DARK = RGBColor(24, 24, 27)            # Code Terminal Dark (#18181B)

# Text & Ink Colors
COLOR_INK_MAIN = RGBColor(15, 23, 42)       # Main Heading Ink (#0F172A)
COLOR_TEXT_BODY = RGBColor(51, 65, 85)      # Body Text (#334155)
COLOR_TEXT_MUTED = RGBColor(100, 116, 139)  # Muted Caption (#64748B)
COLOR_TEXT_ON_DARK = RGBColor(248, 250, 252) # Text on Dark (#F8FAFC)
COLOR_TEXT_ON_DARK_MUTED = RGBColor(148, 163, 184) # Muted on Dark (#94A3B8)
COLOR_CODE_GREEN = RGBColor(74, 222, 128)   # Terminal Code Green (#4ADE80)
COLOR_CODE_CYAN = RGBColor(56, 189, 248)    # Terminal Code Cyan (#38BDF8)
COLOR_CODE_YELLOW = RGBColor(253, 224, 71)  # Terminal Code Yellow (#FDE047)
COLOR_CODE_RED = RGBColor(248, 113, 113)    # Terminal Code Red (#F87171)

# Hairlines & Borders
COLOR_BORDER_SUBTLE = RGBColor(226, 232, 240) # Border Subtle (#E2E8F0)
COLOR_BORDER_DARK = RGBColor(63, 63, 70)      # Border Dark Subtle (#3F3F46)

blank_layout = prs.slide_layouts[6]

def set_slide_background(slide, color):
    """슬라이드 전체 배경색 설정"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_fixed_header(slide, category, title, subtitle):
    """고정 좌표 헤더 렌더링"""
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.73), Inches(0.30))
    tf_c = cat_box.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    p_c = tf_c.paragraphs[0]
    p_c.text = category.upper()
    p_c.font.name = FONT_NAME
    p_c.font.size = Pt(10.5)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.73), Inches(0.45))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_NAME
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_INK_MAIN
    
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.20), Inches(11.73), Inches(0.35))
    tf_s = sub_box.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
    p_s = tf_s.paragraphs[0]
    p_s.text = subtitle
    p_s.font.name = FONT_NAME
    p_s.font.size = Pt(12)
    p_s.font.color.rgb = COLOR_TEXT_MUTED

def add_card(slide, left, top, width, height, title, items, bg_color=COLOR_CARD_WHITE, border_color=COLOR_BORDER_SUBTLE):
    """간결화된 텍스트 카드"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.28), top + Inches(0.24), width - Inches(0.56), height - Inches(0.48))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(15)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_INK_MAIN
    p_title.space_after = Pt(10)
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.name = FONT_NAME
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_BODY
        p.space_after = Pt(6)
    return card

def add_code_card(slide, left, top, width, height, title, code_snippets, header_color=COLOR_CODE_CYAN):
    """다크 터미널/코드 박스"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_DARK
    card.line.color.rgb = COLOR_BORDER_DARK
    card.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.20), width - Inches(0.5), height - Inches(0.40))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = f"💻 {title}"
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = header_color
    p_title.space_after = Pt(6)
    
    for line in code_snippets:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_NAME
        p.font.size = Pt(10)
        if line.startswith("#") or line.startswith("//"):
            p.font.color.rgb = COLOR_TEXT_ON_DARK_MUTED
        elif "503" in line or "429" in line or "오류" in line or "UNAVAILABLE" in line or "❌" in line or "⚠️" in line:
            p.font.color.rgb = COLOR_CODE_RED
        elif "True" in line or "200" in line or "성공" in line or "✅" in line or "정상" in line or "passed" in line:
            p.font.color.rgb = COLOR_CODE_GREEN
        elif "Gate" in line or "Score" in line or "├──" in line or "└──" in line or "자동 전환" in line:
            p.font.color.rgb = COLOR_CODE_YELLOW
        elif "{" in line or "}" in line or ":" in line:
            p.font.color.rgb = RGBColor(228, 228, 231)
        else:
            p.font.color.rgb = COLOR_TEXT_ON_DARK
        p.space_after = Pt(2.5)
    return card

def add_takeaway_strip(slide, left, top, width, height, title, description, accent_color=COLOR_PRIMARY):
    """하단 테이크어웨이 스트립"""
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(241, 245, 249)
    strip.line.color.rgb = accent_color
    strip.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), width - Inches(0.5), height - Inches(0.36))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = f"💡 {title}"
    p.font.name = FONT_NAME
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.space_after = Pt(3)
    
    p2 = tf.add_paragraph()
    p2.text = description
    p2.font.name = FONT_NAME
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_INK_MAIN
    return strip

def add_metric_card(slide, left, top, width, height, label, value, subtext="", accent_color=COLOR_PRIMARY):
    """하단 KPI 지표 카드"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_WHITE
    card.line.color.rgb = COLOR_BORDER_SUBTLE
    card.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_label = tf.paragraphs[0]
    p_label.text = label
    p_label.font.name = FONT_NAME
    p_label.font.size = Pt(10.5)
    p_label.font.color.rgb = COLOR_TEXT_MUTED
    p_label.space_after = Pt(2)
    
    p_val = tf.add_paragraph()
    p_val.text = value
    p_val.font.name = FONT_NAME
    p_val.font.size = Pt(18)
    p_val.font.bold = True
    p_val.font.color.rgb = accent_color
    
    if subtext:
        p_sub = tf.add_paragraph()
        p_sub.text = subtext
        p_sub.font.name = FONT_NAME
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = COLOR_TEXT_BODY
    return card

def add_node(slide, left, top, width, height, text, subtext="", bg_color=COLOR_PRIMARY, text_color=RGBColor(255, 255, 255)):
    """LangGraph 노드 박스"""
    node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    node.fill.solid()
    node.fill.fore_color.rgb = bg_color
    node.line.fill.background()
    
    tf = node.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    
    if subtext:
        p2 = tf.add_paragraph()
        p2.text = subtext
        p2.font.name = FONT_NAME
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = text_color
        p2.alignment = PP_ALIGN.CENTER
    return node

# ==========================================
# Slide 01: 표지 (Cover)
# ==========================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1, COLOR_CANVAS_DARK)

dec = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.6))
dec.fill.solid()
dec.fill.fore_color.rgb = COLOR_EMERALD
dec.line.fill.background()

tb_title = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.3), Inches(3.6))
tf1 = tb_title.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

p_tag = tf1.paragraphs[0]
p_tag.text = "AI AGENT PIPELINE & FULL-STACK VERIFICATION"
p_tag.font.name = FONT_NAME
p_tag.font.size = Pt(12)
p_tag.font.bold = True
p_tag.font.color.rgb = COLOR_EMERALD
p_tag.space_after = Pt(12)

p_main = tf1.add_paragraph()
p_main.text = "🥗 AI 다이어트 코치 에이전트 파이프라인 검증"
p_main.font.name = FONT_NAME
p_main.font.size = Pt(36)
p_main.font.bold = True
p_main.font.color.rgb = COLOR_TEXT_ON_DARK
p_main.space_after = Pt(14)

p_sub = tf1.add_paragraph()
p_sub.text = "식약처 DB · METs 운동 계산 · LangGraph 라우팅 & Self-RAG 3단계 품질 검증"
p_sub.font.name = FONT_NAME
p_sub.font.size = Pt(15)
p_sub.font.color.rgb = RGBColor(203, 213, 225)
p_sub.space_after = Pt(32)

p_info = tf1.add_paragraph()
p_info.text = "발표자 : 메타코드M 라이브 스터디  |  검증 기준 : Diet_Agent_Pipeline_Verification.ipynb"
p_info.font.name = FONT_NAME
p_info.font.size = Pt(11.5)
p_info.font.color.rgb = COLOR_TEXT_ON_DARK_MUTED

# ==========================================
# Slide 02: 목차 (Table of Contents)
# ==========================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2, COLOR_CANVAS_LIGHT)
add_fixed_header(s2, "Table of Contents", "프레젠테이션 목차", "파이프라인 검증 노트북(Diet_Agent_Pipeline_Verification.ipynb) 중심의 체계적 구성")

add_card(s2, Inches(0.8), Inches(1.75), Inches(2.75), Inches(3.4), "Ⅰ. 개요 & MVP 구동", [
    "01. 기획 배경 & 문제 정의",
    "02. 서비스 MVP 구동 요약"
])
add_card(s2, Inches(3.79), Inches(1.75), Inches(2.75), Inches(3.4), "Ⅱ. 핵심 도구 검증", [
    "03. LangGraph 상태 그래프 시각화",
    "04. [Tool 1] 식약처 영양 DB 실측",
    "05. [Tool 2] METs 운동 계산 실측",
    "06. [Tool 3] 영양 백과 RAG 검색"
])
add_card(s2, Inches(6.78), Inches(1.75), Inches(2.75), Inches(3.4), "Ⅲ. 복원력 & Self-RAG", [
    "07. 5대 Flash 모델 자동 폴백",
    "08. HIL 메타데이터 태깅 검증",
    "09. Self-RAG 3단계 품질 게이트",
    "10. SQLite DB & 순 칼로리 집계"
])
add_card(s2, Inches(9.77), Inches(1.75), Inches(2.75), Inches(3.4), "Ⅳ. 총평 & 로드맵", [
    "11. 1~7차시 커리큘럼 연계",
    "12. 스마트 헬스케어 확장 계획"
])

add_takeaway_strip(s2, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45), 
    "핵심 발표 목표 (Key Objective)", 
    "주피터 노트북 환경에서 3대 전문 도구, LangGraph 상태 그래프 시각화, 5대 모델 폴백, Self-RAG 3단계 품질 게이트를 단계별로 실증한 엔지니어링 파이프라인을 전달합니다."
)

# ==========================================
# Slide 03: 기획 배경 및 문제 정의
# ==========================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3, COLOR_CANVAS_LIGHT)
add_fixed_header(s3, "Problem & Solution", "기획 배경 및 해결 과제", "수동 기록의 피로도와 일반 LLM의 환각(Hallucination) 한계를 동시 극복")

add_card(s3, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "⚠️ 기존 다이어트 앱 & LLM의 한계", [
    "수동 입력 피로도: 식사마다 g 수 검색 및 직접 타이핑",
    "치명적 환각 (Hallucination): 임의로 칼로리 수치 왜곡 생성",
    "개인화 결여: 사용자 체형/목표(BMR/TDEE) 미반영",
    "단일 모델 장애: 트래픽 과부하(503/429) 시 서비스 전면 중단"
])

add_card(s3, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "✨ AI 코치의 해결 솔루션", [
    "📸 사진 식단 분석: 멀티모달 비전으로 메뉴 자동 인식",
    "🔍 식약처 표준 DB 도구: Function Calling으로 신뢰성 100%",
    "🛡️ Self-RAG 품질 게이트: 관련성·환각·임상 가드레일 3중 검증",
    "🏃 METs 운동 대사량 연동: 소모 칼로리 산출 및 순 칼로리 관리"
])

add_takeaway_strip(s3, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "핵심 가치 제안 (Value Proposition)",
    "사진 1장 또는 자연어 대화만으로 식약처 표준 영양 데이터를 확인하고, 원클릭으로 DB에 저장하여 실시간 순 칼로리 대시보드와 정기 결산 리포트를 제공합니다."
)

# ==========================================
# Slide 04: [컴팩트] 서비스 MVP 구동 요약
# ==========================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4, COLOR_CANVAS_LIGHT)
add_fixed_header(s4, "MVP Demonstration", "서비스 MVP 구동 요약 (Live Prototype)", "Streamlit 기반 반응형 웹 UI와 4대 대시보드, 텔레그램 정기 결산 연동")

add_card(s4, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🖥️ 클라이언트 주요 화면 구성", [
    "📸 멀티모달 입력: 카메라 촬영 및 식단 이미지 업로드",
    "💬 AI 실시간 대화: 식단/운동/영양 상식 코칭 대화창",
    "🍱 HIL 스마트 카드: AI 분석 결과 확인 후 원클릭 DB 저장",
    "📊 4대 대시보드: 일별 순 칼로리 게이지, 주간/월간/연간 추이",
    "📱 텔레그램 결산: 매월 1일 월간 통계 및 차트 자동 전송"
])

# 우측 4개 핵심 메트릭 카드
add_metric_card(s4, Inches(6.81), Inches(1.75), Inches(2.70), Inches(1.60), "웹 프레임워크", "Streamlit", "반응형 인터랙티브 대시보드", COLOR_PRIMARY)
add_metric_card(s4, Inches(9.83), Inches(1.75), Inches(2.70), Inches(1.60), "인공지능 모델", "Gemini Flash", "초고속 멀티모달 & 도구 바인딩", COLOR_EMERALD)
add_metric_card(s4, Inches(6.81), Inches(3.55), Inches(2.70), Inches(1.60), "데이터베이스", "SQLite + Salt", "SHA-256 개인화 격리 저장", COLOR_PURPLE)
add_metric_card(s4, Inches(9.83), Inches(3.55), Inches(2.70), Inches(1.60), "알림 자동화", "Telegram Bot", "매월 1일 자정 브로드캐스트", COLOR_AMBER)

add_takeaway_strip(s4, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "MVP 핵심 구동 가치 (MVP Value Summary)",
    "복잡한 입력 절차 없이 사진 1장으로 식약처 영양 분석과 운동 소모 칼로리를 자동 반영하여 완성형 풀스택 서비스를 구현했습니다."
)

# ==========================================
# Slide 05: LangGraph 상태 그래프 시각화 (draw_mermaid)
# ==========================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5, COLOR_CANVAS_LIGHT)
add_fixed_header(s5, "LangGraph Workflow", "LangGraph 상태 그래프(StateGraph) 및 시각화 실증", "get_graph().draw_mermaid() 기반의 조건부 라우팅 및 노드 흐름도")

add_card(s5, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🔀 상태 그래프(StateGraph) 구조 설계", [
    "State 정의: `question`, `intent`, `tool_results`, `quality_score`",
    "START ➔ `intent_router`: 질문 의도 4대 조건부 분기",
    "4대 핸들러: `food`, `exercise`, `rag`, `general`",
    "품질 평가: 4대 핸들러 ➔ `self_rag_evaluator` (3단계 검증)",
    "스마트 저장: `human_in_the_loop` ➔ END (DB 저장)"
])

# 우측 실제 Mermaid 시각화 코드 및 출력 터미널 박스
add_code_card(s5, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Live Graph Visualization Output", [
    "# [LangGraph 그래프 시각화 코드 실행]",
    "graph = compiled_agent.get_graph()",
    "mmd = graph.draw_mermaid()",
    "```mermaid",
    "graph TD",
    "  START --> intent_router",
    "  intent_router -.-> food_handler & exercise_handler",
    "  intent_router -.-> rag_handler & general_handler",
    "  food_handler & exercise_handler --> self_rag_evaluator",
    "  rag_handler & general_handler --> self_rag_evaluator",
    "  self_rag_evaluator --> human_in_the_loop --> END",
    "```"
], header_color=COLOR_CODE_CYAN)

add_takeaway_strip(s5, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "상태 그래프 오케스트레이션 (Stateful Orchestration)",
    "LangGraph를 통해 입력 의도에 따라 전문 도구로 분기하고, Self-RAG 품질 검증을 거쳐 스마트 저장으로 이어지는 제어 흐름을 완성했습니다."
)

# ==========================================
# Slide 06: [Tool 1 검증] 식약처 영양 DB 검색
# ==========================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6, COLOR_CANVAS_LIGHT)
add_fixed_header(s6, "Tool 1 Verification", "[Tool 1] 식약처 표준 영양 DB 검색기 실측 검증", "search_food_nutrition() 함수를 통한 5,000+ 공공데이터 실측 수치 반환")

add_card(s6, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🔍 Tool 1 작동 원리 & 검증 포인트", [
    "공공데이터 연동: 식약처 5,000+ 식품 표준 영양 CSV 탑재",
    "키워드 매칭 로직: 사용자가 입력한 대표 음식명 자동 탐색",
    "추출 영양 성분: 칼로리, 탄수화물, 단백질, 지방, 당류, 나트륨",
    "환각율 0%: LLM 추측을 배제하고 CSV 실측치만 정확히 반환"
])

add_code_card(s6, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 3: search_food_nutrition('닭가슴살')", [
    "# [입력]: food_name = '닭가슴살'",
    "# [식약처 CSV DB 실측 반환 JSON]:",
    "{",
    "  '식품명': '콜라겐이첨가된훈제닭가슴살',",
    "  '기준량': '100g',",
    "  '칼로리(kcal)': 135.0,",
    "  '단백질(g)': 26.0,  '지방(g)': 1.5,  '탄수화물(g)': 3.0,",
    "  '나트륨(mg)': 58.0",
    "}",
    "# Status: ✅ 100% 실측 영양 데이터 반환 성공"
], header_color=COLOR_CODE_GREEN)

add_takeaway_strip(s6, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "100% 공공데이터 기반 신뢰성 (Zero-Hallucination Verified)",
    "LLM이 임의로 칼로리를 추측하지 않고 파이썬 함수가 공공데이터 실측치(135kcal, 단백질 26g)를 직접 반환합니다."
)

# ==========================================
# Slide 07: [Tool 2 검증] ACSM METs 운동 소모 칼로리 계산
# ==========================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7, COLOR_CANVAS_LIGHT)
add_fixed_header(s7, "Tool 2 Verification", "[Tool 2] ACSM 표준 METs 운동 소모 계산기 실측 검증", "calculate_exercise_calories() 함수를 통한 과학적 소모 대사량 산출")

add_card(s7, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🏃 Tool 2 작동 원리 & 검증 포인트", [
    "ACSM 공식: `소모 칼로리 = 1.05 × METs × 체중(kg) × 시간(hr)`",
    "20+ 운동 계수 DB: 러닝(8.5), 웨이트(5.5), 수영(7.5), 사이클(7.0)",
    "체중 연동: 사용자별 체중(kg) 자동 반영 맞춤 계산",
    "순 칼로리(Net Calories) 연동: 섭취량 - 운동 소모량"
])

add_code_card(s7, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 4: calculate_exercise_calories('러닝')", [
    "# [입력]: exercise_name='러닝', duration=30분, weight=70kg",
    "# [ACSM 공식 계산 실측 반환 JSON]:",
    "{",
    "  '운동명': '러닝',",
    "  '적용체중(kg)': 70.0,",
    "  'METs계수': 8.5,",
    "  '소모칼로리(kcal)': 312.4,",
    "  '설명': '70.0kg 기준 러닝 30분 수행 시 약 312.4 kcal 소모'",
    "}"
], header_color=COLOR_CODE_GREEN)

add_takeaway_strip(s7, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "과학적 대사량 관리 (Scientific Energy Expenditure)",
    "미국 스포츠의학회(ACSM) 공식에 체중과 운동 시간을 연동하여 정밀한 소모 칼로리를 산출합니다."
)

# ==========================================
# Slide 08: [Tool 3 검증] 영양 백과 RAG 지식 검색
# ==========================================
s8 = prs.slides.add_slide(blank_layout)
set_slide_background(s8, COLOR_CANVAS_LIGHT)
add_fixed_header(s8, "Tool 3 Verification", "[Tool 3] 다이어트 & 영양 백과 RAG 검색기 실측 검증", "search_nutrition_knowledge() 함수를 통한 임상 영양 전문 가이드 제공")

add_card(s8, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "📚 Tool 3 작동 원리 & 검증 포인트", [
    "영양 백과 RAG: 전문 임상 영양 지식 베이스 구축",
    "혈당 스파이크 방지: [식이섬유 ➔ 단백질 ➔ 탄수화물] 식사순서",
    "정체기 극복 가이드: 대사 적응 극복을 위한 리피드(Refeed) 전략",
    "단백질 흡수 타이밍: 운동 직후 근합성을 위한 단백질 권장량"
])

add_code_card(s8, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 5: search_nutrition_knowledge('혈당')", [
    "# [입력]: query = '혈당 스파이크 방지'",
    "# [영양 RAG 지식 검색 실측 반환 JSON]:",
    "{",
    "  '주제': '혈당 스파이크 방지 및 식사 순서',",
    "  '핵심원리': '급격한 인슐린 분비는 체지방 축적 유발',",
    "  '실천가이드': '1. 식이섬유 -> 2. 단백질/지방 -> 3. 탄수화물',",
    "  '추천행동': '식후 15분 가벼운 산책 권장'",
    "}"
], header_color=COLOR_CODE_GREEN)

add_takeaway_strip(s8, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "임상 영양 가이드 (Clinical Nutrition RAG)",
    "단순 수치 계산을 넘어 혈당 관리 및 정체기 극복을 위한 전문 임상 영양 가이드를 과학적으로 제공합니다."
)

# ==========================================
# Slide 09: [무중단 복원력] 5대 Flash 모델 폴백 실증
# ==========================================
s9 = prs.slides.add_slide(blank_layout)
set_slide_background(s9, COLOR_CANVAS_LIGHT)
add_fixed_header(s9, "Fault Tolerance", "5대 Flash 모델 무중단 폴백(Fallback) 실증", "503 과부하 및 429 속도제한 감지 시 0.5초 이내 예비 Flash 모델 자동 전환")

add_card(s9, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🛡️ 다중 모델 폴백 아키텍처", [
    "1순위: `gemini-3.6-flash` (기본 호출 & 3대 Tool 바인딩)",
    "2순위: `gemini-3.7-flash` (차세대 플래시 1차 예비)",
    "3순위: `gemini-3.5-flash` (고가용성 플래시 2차 예비)",
    "4순위: `gemini-flash-latest` (최신 안정화 에일리어스)",
    "5순위: `gemini-2.5-flash-lite` (초경량 최저 토큰 소모)"
])

add_code_card(s9, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 6: Self-Healing Fallback Test", [
    "# [1. 음식 사진 식단 분석 요청 전송]",
    "Attempt 1: Call [gemini-3.6-flash] with food image...",
    "⚠️ [gemini-3.6-flash] 503 UNAVAILABLE (Demand Spikes)",
    "",
    "# [2. 시스템 자동 감지 & 0.5s 내 예비 모델 전환]",
    "🔄 Auto-Switching: [gemini-3.7-flash]",
    "Attempt 2: Call [gemini-3.7-flash] with tools & session...",
    "",
    "# [3. 무중단 정상 응답 복구 완료]",
    "✅ [gemini-3.7-flash] 200 OK (분석 성공 & 메타데이터 추출)"
], header_color=COLOR_ROSE)

add_takeaway_strip(s9, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "고가용성 복원력 (High Availability Resilience)",
    "특정 모델의 트래픽 폭주 시에도 0.5초 이내에 예비 모델로 자동 전환되어 중단 없는 서비스를 유지합니다."
)

# ==========================================
# Slide 10: [멀티모달 & HIL] 식단 분석 및 메타데이터 태깅
# ==========================================
s10 = prs.slides.add_slide(blank_layout)
set_slide_background(s10, COLOR_CANVAS_LIGHT)
add_fixed_header(s10, "Multimodal & HIL", "멀티모달 식단 분석 & 구조화 메타데이터 태깅", "식단 분석 결과에서 MEAL_DATA JSON 태그 자동 생성 및 파싱 검증")

add_card(s10, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🏷️ Human-in-the-Loop 메타데이터 태깅 원리", [
    "3단계 프롬프트: ① 영양소 요약 ➔ ② 목표 대비 진단 ➔ ③ 메뉴 제안",
    "메타데이터 생성: `<!-- MEAL_DATA: {...} -->` 한 줄 태깅",
    "정규식 파서: `parse_agent_metadata()`로 텍스트와 JSON 분리",
    "스마트 컨펌 카드: UI 하단에 즉시 렌더링되어 원클릭 DB 저장"
])

add_code_card(s10, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 7: parse_agent_metadata() Output", [
    "# [AI 코치 클린 답변 텍스트]:",
    "오늘 점심으로 드신 닭가슴살(135kcal)과 사과(85kcal)는",
    "총 220kcal이며 단백질 26g을 충분히 보충하셨습니다...",
    "",
    "# [추출된 식단 메타데이터 (MEAL_DATA)]:",
    "{",
    "  'food_name': '닭가슴살과 사과',",
    "  'calories': 220,  'carbs': 25,  'protein': 26,",
    "  'fat': 2,  'sodium': 60,  'meal_type': '점심'",
    "}"
], header_color=COLOR_CODE_CYAN)

add_takeaway_strip(s10, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "정형 데이터 자동 추출 (Structured Data Extraction)",
    "자연어 답변에서 정형 JSON 데이터를 완벽히 분리하여 사용자의 원클릭 DB 저장을 지원합니다."
)

# ==========================================
# Slide 11: [Self-RAG 품질 검증] LLM-as-a-Judge 3단계
# ==========================================
s11 = prs.slides.add_slide(blank_layout)
set_slide_background(s11, COLOR_CANVAS_LIGHT)
add_fixed_header(s11, "Self-RAG Quality Gate", "Self-RAG 3단계 품질 게이트 (LLM-as-a-Judge) 실증", "관련성(Relevance) ➔ 환각 검출(Grounding) ➔ 임상 안전성(Safety) 3중 검증")

add_card(s11, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🛡️ 3단계 이진 품질 평가 구조", [
    "Gate 1. 영양 관련성 검증 (Relevance Check):\n  • 사용자 식단 ↔ DB 조회 식품 일치 여부 판정 (Yes/No)\n  • 불일치 시 쿼리 자동 재작성 (CRAG 보정)",
    "Gate 2. 환각 수치 검출 (Hallucination Grounding):\n  • 칼로리/탄단지 수치와 DB 실측치 100% 일치 검증",
    "Gate 3. 임상 안전 가드레일 (Clinical Safety):\n  • 초저열량 경고 및 의학적 면책 안내 자동 삽입"
])

add_code_card(s11, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 9: evaluate_self_rag_quality()", [
    "# [Self-RAG 3단계 품질 게이트 실행 결과]:",
    "{",
    "  'Gate_1_Relevance': 'yes',",
    "  'Gate_2_Grounding_Zero_Hallucination': True,",
    "  'Gate_3_Clinical_Safety': 'safe',",
    "  'ALL_GATES_PASSED': True",
    "}",
    "# Status: ✅ 3단계 품질 게이트 ALL PASS (최종 승인 반환)"
], header_color=COLOR_CODE_GREEN)

add_takeaway_strip(s11, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "철저한 환각 통제 (Zero-Hallucination Verified)",
    "LLM-as-a-Judge 평가를 통해 관련성, 사실 근거성, 임상 안전성을 3중 검증하여 안전한 답변만 제공합니다."
)

# ==========================================
# Slide 12: [DB & 순 칼로리] SQLite 연동 실측 쿼리
# ==========================================
s12 = prs.slides.add_slide(blank_layout)
set_slide_background(s12, COLOR_CANVAS_LIGHT)
add_fixed_header(s12, "Database & Net Calories", "SQLite DB 연동 및 순 칼로리 집계 쿼리 실측", "get_daily_summary() 함수를 통한 섭취량 - 소모량 실시간 집계")

add_card(s12, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "💾 SQLite 데이터베이스 & 순 칼로리 연산", [
    "테이블 구조: `users`, `meal_records`, `exercise_records`",
    "보안 암호화: SHA-256 + 32바이트 Salt 난수 해싱",
    "순 칼로리 공식: `Net Calories = 총 섭취 칼로리 - 총 운동 소모량`",
    "대시보드 실시간 반영: 일별 순 칼로리 게이지 및 도넛 차트"
])

add_code_card(s12, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "Notebook Cell 10: get_daily_summary(user_id=1)", [
    "# [SQLite 일별 종합 집계 실측 쿼리 결과]:",
    "{",
    "  'date': '2026-08-29',",
    "  'total_cal': 520.0,       # 섭취 칼로리",
    "  'total_burned': 312.4,    # 운동 소모 칼로리",
    "  'net_cal': 207.6,         # ✨ 순 칼로리 (520 - 312.4)",
    "  'total_protein': 36.4,    # 섭취 단백질(g)",
    "  'records_count': 2,       'exercise_count': 1",
    "}"
], header_color=COLOR_CODE_YELLOW)

# 하단 4개 지표 카드
add_metric_card(s12, Inches(0.8), Inches(5.35), Inches(2.70), Inches(1.45), "Daily Intake", "520.0 kcal", "일일 목표 2,000 kcal 대비", COLOR_PRIMARY)
add_metric_card(s12, Inches(3.81), Inches(5.35), Inches(2.70), Inches(1.45), "Calories Burned", "-312.4 kcal", "🔥 러닝 30분 소모 실측", COLOR_EMERALD)
add_metric_card(s12, Inches(6.82), Inches(5.35), Inches(2.70), Inches(1.45), "Net Calories", "207.6 kcal", "✨ 순 섭취 칼로리", COLOR_PURPLE)
add_metric_card(s12, Inches(9.83), Inches(5.35), Inches(2.70), Inches(1.45), "Protein Intake", "36.4 g", "일일 목표 100g 대비", COLOR_AMBER)

# ==========================================
# Slide 13: 강의 커리큘럼 연계 및 기술적 의의
# ==========================================
s13 = prs.slides.add_slide(blank_layout)
set_slide_background(s13, COLOR_CANVAS_LIGHT)
add_fixed_header(s13, "Course Mapping", "라이브스터디(1~7차시) 커리큘럼 연계", "강의 핵심 이론 및 프레임워크를 실전 풀스택 프로덕션 서비스로 완성")

add_card(s13, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "📚 차시별 이론 접목 내역", [
    "1차시: 페르소나 및 3단계 응답 구조 프롬프트 엔지니어링",
    "2~4차시: 식약처 CSV 데이터 정제 및 키워드 매칭 로직",
    "5차시: Function Calling 도구 바인딩 & ReAct 에이전트",
    "6차시: 세션 히스토리 메모리 & 멀티모달 사진 처리",
    "7차시: Self-RAG 3단계 품질 게이트 & LangGraph 조건부 라우팅"
])

add_card(s13, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "🏆 프로젝트의 기술적 의의", [
    "풀스택 완성: LLM 단독 실행이 아닌 DB, UI, 메신저 결합",
    "환각 통제: 공공데이터 도구 강제로 의료/영양 신뢰성 확보",
    "무중단 복원력: 5대 Flash 모델 폴백으로 503/429 장애 극복",
    "검증 완결성: 주피터 노트북을 통한 단계별 실측 완료"
])

add_takeaway_strip(s13, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "학습 성과 요약 (Learning Takeaway)",
    "프롬프트부터 RAG, Tool Use, LangGraph, Self-RAG까지 전 과정을 실전 서비스로 구현하여 엔지니어링 역량을 극대화했습니다."
)

# ==========================================
# Slide 14: 배포 성과 및 향후 발전 로드맵
# ==========================================
s14 = prs.slides.add_slide(blank_layout)
set_slide_background(s14, COLOR_CANVAS_LIGHT)
add_fixed_header(s14, "Summary & Future Work", "서비스 배포 성과 & 발전 로드맵", "글로벌 배포 완료 및 스마트 헬스케어 생태계로의 확장 계획")

add_card(s14, Inches(0.8), Inches(1.75), Inches(5.72), Inches(3.4), "🚀 배포 성과 & 서비스 현황", [
    "GitHub 저장소: rye6837-web/impossible_to_get_lost_my_weight",
    "Streamlit Community Cloud: 모바일/PC 반응형 웹 서비스 운영 중",
    "핵심 가치: 사진 1장 식단 분석, 운동 칼로리 차감, 결산 자동화",
    "안정성: 다중 모델 폴백 체계로 무중단 24/7 서비스 유지"
])

add_card(s14, Inches(6.81), Inches(1.75), Inches(5.72), Inches(3.4), "🔮 향후 확장 로드맵 (Roadmap)", [
    "스마트워치 연동: 애플워치·갤럭시워치 활동 칼로리 실시간 동기화",
    "AI 맞춤 식단 플래너: 일주일 식단표 자동 생성 및 밀키트 연계",
    "연속 혈당 측정기(CGM) 접목: 혈당 반응 기반 실시간 코칭 추가"
])

add_takeaway_strip(s14, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.45),
    "감사 인사 및 Q&A",
    "경청해 주셔서 감사합니다. 질문 및 피드백을 환영합니다!"
)

# 저장
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AI_Diet_Coach_Presentation.pptx")
prs.save(output_path)
print(f"✅ 검증 노트북 기반 14장 PPT 생성 완료: {output_path}")
